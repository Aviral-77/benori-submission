import csv
import io
import json
from datetime import datetime, timezone

RAW_FIELDS = [
    "rank_score", "relevance", "credibility", "is_relevant", "deal_type",
    "deal_value", "acquirer", "target", "title", "publisher", "source_domain",
    "credibility_tier", "is_press_release", "corroboration_count",
    "cluster_size", "published", "url", "matched_deal_terms",
    "matched_fmcg_terms", "duplicate_urls",
]


def _flat(value):
    if isinstance(value, (list, tuple)):
        return "; ".join(str(v) for v in value)
    return "" if value is None else str(value)


def to_csv(articles):
    buf = io.StringIO()
    writer = csv.DictWriter(buf, fieldnames=RAW_FIELDS, extrasaction="ignore")
    writer.writeheader()
    for art in articles:
        writer.writerow({k: _flat(art.get(k)) for k in RAW_FIELDS})
    return buf.getvalue().encode("utf-8")


def to_json(articles):
    cleaned = [{k: art.get(k) for k in RAW_FIELDS} for art in articles]
    return json.dumps(cleaned, indent=2, ensure_ascii=False).encode("utf-8")


def to_excel(newsletter, articles):
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment, PatternFill
    from openpyxl.utils import get_column_letter

    wb = Workbook()
    accent = PatternFill("solid", fgColor="1F4E79")
    head_font = Font(bold=True, color="FFFFFF")
    title_font = Font(bold=True, size=14, color="1F4E79")

    ws = wb.active
    ws.title = "Newsletter"
    ws["A1"] = newsletter["title"]
    ws["A1"].font = title_font
    ws["A2"] = newsletter["subtitle"]
    ws["A2"].font = Font(italic=True, color="555555")

    headers = ["#", "Headline", "Deal type", "Value", "Parties",
               "Summary", "Source", "Credibility", "Published", "Link"]
    row0 = 4
    for c, h in enumerate(headers, start=1):
        cell = ws.cell(row=row0, column=c, value=h)
        cell.fill, cell.font = accent, head_font
    for i, item in enumerate(newsletter["lead_deals"], start=1):
        r = row0 + i
        parties = " -> ".join(p for p in [item.get("acquirer"), item.get("target")] if p)
        ws.cell(row=r, column=1, value=i)
        ws.cell(row=r, column=2, value=item["title"])
        ws.cell(row=r, column=3, value=item.get("deal_type"))
        ws.cell(row=r, column=4, value=item.get("deal_value") or "-")
        ws.cell(row=r, column=5, value=parties or "-")
        ws.cell(row=r, column=6, value=item["summary"])
        ws.cell(row=r, column=7, value=item.get("publisher"))
        ws.cell(row=r, column=8, value=f'{item.get("credibility")} ({item.get("credibility_tier")})')
        ws.cell(row=r, column=9, value=item.get("published", "")[:10])
        ws.cell(row=r, column=10, value=item.get("url"))
    widths = [4, 46, 14, 12, 26, 60, 18, 22, 12, 40]
    for c, w in enumerate(widths, start=1):
        ws.column_dimensions[get_column_letter(c)].width = w
    for row in ws.iter_rows(min_row=row0 + 1, max_row=row0 + len(newsletter["lead_deals"])):
        for cell in row:
            cell.alignment = Alignment(vertical="top", wrap_text=True)

    ws2 = wb.create_sheet("Raw data")
    for c, h in enumerate(RAW_FIELDS, start=1):
        cell = ws2.cell(row=1, column=c, value=h)
        cell.fill, cell.font = accent, head_font
    for i, art in enumerate(articles, start=2):
        for c, k in enumerate(RAW_FIELDS, start=1):
            ws2.cell(row=i, column=c, value=_flat(art.get(k)))
    for c in range(1, len(RAW_FIELDS) + 1):
        ws2.column_dimensions[get_column_letter(c)].width = 22

    ws3 = wb.create_sheet("Methodology")
    ws3.column_dimensions["A"].width = 110
    ws3["A1"] = "Pipeline & assumptions"
    ws3["A1"].font = title_font
    for i, line in enumerate(newsletter["methodology"], start=3):
        ws3.cell(row=i, column=1, value=line).alignment = Alignment(wrap_text=True)

    out = io.BytesIO()
    wb.save(out)
    return out.getvalue()


def to_word(newsletter):
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()
    accent = RGBColor(0x1F, 0x4E, 0x79)

    h = doc.add_heading(newsletter["title"], level=0)
    h.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub = doc.add_paragraph(newsletter["subtitle"])
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub.runs[0].italic = True

    intro = doc.add_paragraph(newsletter["intro"])
    intro.runs[0].font.size = Pt(11)

    doc.add_heading("Lead deals", level=1)
    for i, item in enumerate(newsletter["lead_deals"], start=1):
        p = doc.add_paragraph()
        run = p.add_run(f"{i}. {item['title']}")
        run.bold = True
        run.font.color.rgb = accent
        run.font.size = Pt(12)

        meta_bits = [item.get("deal_type")]
        if item.get("deal_value"):
            meta_bits.append(item["deal_value"])
        parties = " -> ".join(x for x in [item.get("acquirer"), item.get("target")] if x)
        if parties:
            meta_bits.append(parties)
        meta = doc.add_paragraph(" • ".join(b for b in meta_bits if b))
        meta.runs[0].italic = True
        meta.runs[0].font.size = Pt(9)

        doc.add_paragraph(item["summary"])

        src = doc.add_paragraph()
        s = src.add_run(
            f"Source: {item.get('publisher')} - credibility "
            f"{item.get('credibility')}/100 ({item.get('credibility_label')})"
            + (f" • corroborated by {item.get('corroboration_count')} other outlet(s)"
               if item.get("corroboration_count") else "")
        )
        s.font.size = Pt(8)
        s.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        link = doc.add_paragraph()
        lr = link.add_run(item.get("url", ""))
        lr.font.size = Pt(8)
        lr.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)

    if newsletter["brief_mentions"]:
        doc.add_heading("Also in the news", level=1)
        for item in newsletter["brief_mentions"]:
            b = doc.add_paragraph(style="List Bullet")
            br = b.add_run(f"{item['title']} - {item.get('publisher')}")
            br.font.size = Pt(10)

    doc.add_heading("Methodology & assumptions", level=1)
    for line in newsletter["methodology"]:
        mp = doc.add_paragraph(line)
        mp.runs[0].font.size = Pt(8)
        mp.runs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    out = io.BytesIO()
    doc.save(out)
    return out.getvalue()


def to_pptx(newsletter):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    accent = RGBColor(0x1F, 0x4E, 0x79)
    grey = RGBColor(0x66, 0x66, 0x66)
    blank = prs.slide_layouts[6]

    def textbox(slide, left, top, width, height):
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tb.text_frame.word_wrap = True
        return tb.text_frame

    s = prs.slides.add_slide(blank)
    tf = textbox(s, 0.8, 2.4, 11.7, 2.2)
    tf.text = newsletter["title"]
    tf.paragraphs[0].font.size = Pt(40)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = accent
    p = tf.add_paragraph()
    p.text = newsletter["subtitle"]
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = grey

    s = prs.slides.add_slide(blank)
    tf = textbox(s, 0.8, 0.5, 11.7, 1.0)
    tf.text = "This edition at a glance"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = accent
    body = textbox(s, 0.8, 1.7, 11.7, 5.0)
    body.text = newsletter["intro"]
    body.paragraphs[0].font.size = Pt(16)
    for stat in newsletter["stats_lines"]:
        pp = body.add_paragraph()
        pp.text = "•  " + stat
        pp.font.size = Pt(14)
        pp.font.color.rgb = grey

    for i, item in enumerate(newsletter["lead_deals"], start=1):
        s = prs.slides.add_slide(blank)
        head = textbox(s, 0.8, 0.5, 11.7, 1.4)
        head.text = f"{i}. {item['title']}"
        head.paragraphs[0].font.size = Pt(24)
        head.paragraphs[0].font.bold = True
        head.paragraphs[0].font.color.rgb = accent

        meta_bits = [item.get("deal_type")]
        if item.get("deal_value"):
            meta_bits.append(item["deal_value"])
        parties = " -> ".join(x for x in [item.get("acquirer"), item.get("target")] if x)
        if parties:
            meta_bits.append(parties)
        meta = textbox(s, 0.8, 1.9, 11.7, 0.6)
        meta.text = " • ".join(b for b in meta_bits if b)
        meta.paragraphs[0].font.size = Pt(14)
        meta.paragraphs[0].font.italic = True
        meta.paragraphs[0].font.color.rgb = grey

        body = textbox(s, 0.8, 2.7, 11.7, 3.0)
        body.text = item["summary"]
        body.paragraphs[0].font.size = Pt(16)

        foot = textbox(s, 0.8, 6.4, 11.7, 0.8)
        foot.text = (
            f"Source: {item.get('publisher')}  |  credibility {item.get('credibility')}/100  "
            f"({item.get('credibility_tier')})  |  {item.get('url','')}"
        )
        foot.paragraphs[0].font.size = Pt(10)
        foot.paragraphs[0].font.color.rgb = grey

    s = prs.slides.add_slide(blank)
    tf = textbox(s, 0.8, 0.5, 11.7, 1.0)
    tf.text = "Methodology & assumptions"
    tf.paragraphs[0].font.size = Pt(26)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = accent
    body = textbox(s, 0.8, 1.6, 11.7, 5.4)
    body.text = ""
    for line in newsletter["methodology"]:
        pp = body.add_paragraph()
        pp.text = "•  " + line
        pp.font.size = Pt(12)
        pp.font.color.rgb = grey

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def export_all(newsletter, articles):
    stamp = datetime.now(timezone.utc).strftime("%Y%m%d")
    return {
        f"fmcg_deals_raw_{stamp}.csv": to_csv(articles),
        f"fmcg_deals_raw_{stamp}.json": to_json(articles),
        f"fmcg_newsletter_{stamp}.xlsx": to_excel(newsletter, articles),
        f"fmcg_newsletter_{stamp}.docx": to_word(newsletter),
        f"fmcg_newsletter_{stamp}.pptx": to_pptx(newsletter),
    }
